"""
Document Generator for OpenClaw
Creates PDF, Word, Excel, PowerPoint, and text documents.
Usage: python doc_generator.py <format> <title> <content_file>
"""
import sys
import os
import json
import time

def create_pdf(title, content, output_path):
    """Create a professional PDF document."""
    from fpdf import FPDF
    
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    
    # Title
    pdf.set_font("Helvetica", "B", 24)
    pdf.cell(0, 15, title, ln=True, align="C")
    pdf.ln(5)
    
    # Date
    pdf.set_font("Helvetica", "I", 10)
    pdf.cell(0, 8, time.strftime("Generated: %Y-%m-%d %H:%M"), ln=True, align="C")
    pdf.ln(10)
    
    # Content
    pdf.set_font("Helvetica", "", 11)
    for line in content.split("\n"):
        if line.startswith("# "):
            pdf.set_font("Helvetica", "B", 18)
            pdf.cell(0, 12, line[2:], ln=True)
            pdf.set_font("Helvetica", "", 11)
        elif line.startswith("## "):
            pdf.set_font("Helvetica", "B", 14)
            pdf.cell(0, 10, line[3:], ln=True)
            pdf.set_font("Helvetica", "", 11)
        elif line.startswith("- "):
            pdf.cell(10)
            pdf.cell(0, 7, chr(8226) + " " + line[2:], ln=True)
        elif line.strip():
            pdf.multi_cell(0, 7, line)
        else:
            pdf.ln(3)
    
    pdf.output(output_path)
    return f"PDF created: {output_path}"

def create_word(title, content, output_path):
    """Create a professional Word document."""
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    
    doc = Document()
    
    # Title
    title_para = doc.add_heading(title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Date
    date_para = doc.add_paragraph(time.strftime("Generated: %Y-%m-%d %H:%M"))
    date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph("")  # spacer
    
    # Content
    for line in content.split("\n"):
        if line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("- "):
            doc.add_paragraph(line[2:], style='List Bullet')
        elif line.startswith("1. ") or line.startswith("2. ") or line.startswith("3. "):
            doc.add_paragraph(line[3:], style='List Number')
        elif line.strip():
            doc.add_paragraph(line)
    
    doc.save(output_path)
    return f"Word document created: {output_path}"

def create_excel(title, content, output_path):
    """Create an Excel spreadsheet from structured data."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill
    
    wb = Workbook()
    ws = wb.active
    ws.title = title[:31]  # Excel limit
    
    # Title row
    ws['A1'] = title
    ws['A1'].font = Font(size=16, bold=True)
    ws.merge_cells('A1:E1')
    
    # Date
    ws['A2'] = time.strftime("Generated: %Y-%m-%d %H:%M")
    ws['A2'].font = Font(italic=True, size=9)
    
    # Content - try to parse as table data
    row = 4
    for line in content.split("\n"):
        if "|" in line and "---" not in line:
            cells = [c.strip() for c in line.split("|") if c.strip()]
            for col, cell in enumerate(cells, 1):
                ws.cell(row=row, column=col, value=cell)
                if row == 4:  # header row
                    ws.cell(row=row, column=col).font = Font(bold=True)
                    ws.cell(row=row, column=col).fill = PatternFill("solid", fgColor="4472C4")
                    ws.cell(row=row, column=col).font = Font(bold=True, color="FFFFFF")
            row += 1
        elif line.startswith("- "):
            ws.cell(row=row, column=1, value=line[2:])
            row += 1
        elif line.strip() and not line.startswith("#"):
            ws.cell(row=row, column=1, value=line)
            row += 1
    
    # Auto-width
    for col in ws.columns:
        max_len = max((len(str(cell.value or "")) for cell in col), default=10)
        ws.column_dimensions[col[0].column_letter].width = min(max_len + 2, 50)
    
    wb.save(output_path)
    return f"Excel created: {output_path}"

def create_pptx(title, content, output_path):
    """Create a PowerPoint presentation."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    
    prs = Presentation()
    
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = time.strftime("Generated: %Y-%m-%d %H:%M")
    
    # Content slides - split by headers
    current_title = ""
    current_bullets = []
    
    for line in content.split("\n"):
        if line.startswith("# ") or line.startswith("## "):
            # Save previous slide
            if current_title and current_bullets:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = current_title
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()
                for i, bullet in enumerate(current_bullets):
                    if i == 0:
                        tf.text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
            current_title = line.lstrip("# ").strip()
            current_bullets = []
        elif line.startswith("- "):
            current_bullets.append(line[2:])
        elif line.strip():
            current_bullets.append(line.strip())
    
    # Last slide
    if current_title and current_bullets:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = current_title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for i, bullet in enumerate(current_bullets):
            if i == 0:
                tf.text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
    
    prs.save(output_path)
    return f"PowerPoint created: {output_path}"

def create_txt(title, content, output_path):
    """Create a plain text file."""
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(f"{'=' * 60}\n")
        f.write(f"  {title}\n")
        f.write(f"  {time.strftime('%Y-%m-%d %H:%M')}\n")
        f.write(f"{'=' * 60}\n\n")
        f.write(content)
    return f"Text file created: {output_path}"

# CLI
if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python doc_generator.py <format> <title> <content_or_file>")
        print("Formats: pdf, word, excel, pptx, txt")
        print("Content: direct text OR path to a .md/.txt file")
        sys.exit(1)
    
    fmt = sys.argv[1].lower()
    title = sys.argv[2]
    content_arg = " ".join(sys.argv[3:])
    
    # Check if content is a file path
    if os.path.isfile(content_arg):
        with open(content_arg, "r", encoding="utf-8") as f:
            content = f.read()
    else:
        content = content_arg
    
    # Output path
    ext_map = {"pdf": ".pdf", "word": ".docx", "excel": ".xlsx", "pptx": ".pptx", "txt": ".txt"}
    ext = ext_map.get(fmt, ".txt")
    safe_title = "".join(c if c.isalnum() or c in " -_" else "" for c in title).strip().replace(" ", "_")
    output_path = os.path.join(os.path.dirname(__file__), f"{safe_title}{ext}")
    
    generators = {
        "pdf": create_pdf,
        "word": create_word,
        "excel": create_excel,
        "pptx": create_pptx,
        "txt": create_txt,
    }
    
    if fmt in generators:
        result = generators[fmt](title, content, output_path)
        print(result)
    else:
        print(f"Unknown format: {fmt}. Use: pdf, word, excel, pptx, txt")
        sys.exit(1)
